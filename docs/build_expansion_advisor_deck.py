"""Build a PowerPoint overview of the Expansion Advisor.

Run:
    python docs/build_expansion_advisor_deck.py
Outputs: docs/expansion_advisor_overview.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "expansion_advisor_overview.pptx"

# Theme palette
NAVY = RGBColor(0x0F, 0x2C, 0x3F)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
SAND = RGBColor(0xE8, 0xDD, 0xC2)
INK = RGBColor(0x1A, 0x1A, 0x1A)
SUBTLE = RGBColor(0x55, 0x6B, 0x7A)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT = RGBColor(0xC9, 0x6F, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xB7, 0x1C, 0x1C)
GREY = RGBColor(0x9A, 0xA5, 0xAE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=14, color=INK,
                bullet_color=TEAL, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet marker
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r1.font.name = "Calibri"
        # body
        if isinstance(item, tuple):
            head, tail = item
            r2 = p.add_run()
            r2.text = head
            r2.font.size = Pt(size)
            r2.font.bold = True
            r2.font.color.rgb = color
            r2.font.name = "Calibri"
            if tail:
                r3 = p.add_run()
                r3.text = " — " + tail
                r3.font.size = Pt(size)
                r3.font.color.rgb = color
                r3.font.name = "Calibri"
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = "Calibri"
    return tb


def slide_header(slide, eyebrow, title, page=None):
    # Top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.55), NAVY)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(8), Inches(0.32),
             "OAKTREE ATLAS  ·  EXPANSION ADVISOR", size=11, bold=True,
             color=SAND)
    if page is not None:
        add_text(slide, SLIDE_W - Inches(1.5), Inches(0.12), Inches(1.0),
                 Inches(0.32), page, size=11, bold=True, color=SAND,
                 align=PP_ALIGN.RIGHT)
    # Eyebrow + title
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.35),
             eyebrow, size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    # underline
    add_rect(slide, Inches(0.5), Inches(1.65), Inches(1.2), Inches(0.05), TEAL)


def add_card(slide, x, y, w, h, title, lines, *, head_color=TEAL):
    add_rect(slide, x, y, w, h, LIGHT)
    add_rect(slide, x, y, w, Inches(0.45), head_color)
    add_text(slide, x + Inches(0.15), y + Inches(0.08), w - Inches(0.2),
             Inches(0.32), title, size=13, bold=True, color=WHITE)
    inner_y = y + Inches(0.55)
    add_bullets(slide, x + Inches(0.15), inner_y, w - Inches(0.3),
                h - Inches(0.6), lines, size=11, line_spacing=1.1)


def add_footer(slide, text):
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4), Inches(12.3),
             Inches(0.3), text, size=10, color=SUBTLE)


# ────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(1.9), TEAL)
add_rect(s, Inches(0.7), Inches(2.0), Inches(0.18), Inches(2.5), ACCENT)
add_text(s, Inches(1.0), Inches(1.8), Inches(11), Inches(0.5),
         "OAKTREE ATLAS", size=14, bold=True, color=SAND)
add_text(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.2),
         "Expansion Advisor", size=54, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(3.4), Inches(11.3), Inches(0.6),
         "How it works: data, scoring, gates, and the UI",
         size=24, color=SAND)
add_text(s, Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.5),
         "Riyadh-first restaurant & retail location intelligence",
         size=16, color=SAND)
add_text(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.5),
         "A walk-through for analysts and operators",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.4),
         "Covers data flow · gates · scoring · Decision Memory · UI",
         size=13, color=SAND)

# ────────────────────────────────────────────────────────────────────────────
# Slide 2 — What is the Expansion Advisor
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "01 · OVERVIEW", "What the Expansion Advisor does", "2")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.6),
         "An AI-assisted location-recommendation engine for F&B brands "
         "expanding in Riyadh.", size=18, color=INK)

add_card(s, Inches(0.5), Inches(2.7), Inches(4.0), Inches(4.3),
         "INPUT", [
             ("Brand profile", "name, category, price tier"),
             ("Service model", "QSR, dine-in, delivery-first, café"),
             ("Area & districts", "min/max m², target districts, bbox"),
             ("Existing branches", "for cannibalization analysis"),
             ("Expansion goal", "growth shape and tolerance"),
         ])
add_card(s, Inches(4.7), Inches(2.7), Inches(4.0), Inches(4.3),
         "PROCESS", [
             ("Candidate pool", "stratified by district"),
             ("Context enrichment", "roads, parking, delivery, comps"),
             ("Gates", "zoning, area, access, economics, ..."),
             ("Score 6 dimensions", "weighted composite 0–100"),
             ("LLM rerank + memo", "structured, cached, pre-warmed"),
         ])
add_card(s, Inches(8.9), Inches(2.7), Inches(4.0), Inches(4.3),
         "OUTPUT", [
             ("Ranked candidates", "with score breakdown"),
             ("Decision Memory", "narrative + diagnostics"),
             ("Verdict", "GO / CAUTION / RECONSIDER"),
             ("Comparable competitors", "and risk list"),
             ("Compare & Report views", "shortlist sharing"),
         ])
add_footer(s, "Service version: expansion_advisor_v7  ·  Riyadh only  ·  parcel source: listings_only (Aqar, Wasalt, Bayut)")

# ────────────────────────────────────────────────────────────────────────────
# Slide 3 — End-to-end flow
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "02 · FLOW", "End-to-end pipeline at a glance", "3")

steps = [
    ("Search", "brand, area,\ndistricts"),
    ("Pool", "parcels +\nlistings"),
    ("Enrich", "roads · parking ·\ndelivery · comps"),
    ("Gates", "structural &\noperational"),
    ("Score", "6 weighted\ndimensions"),
    ("Rerank", "LLM shortlist\n(optional)"),
    ("Memo", "structured\ndecision memory"),
    ("UI", "results +\ndiagnostics"),
]
n = len(steps)
total_w = Inches(12.4)
gap = Inches(0.10)
box_w = (total_w - gap * (n - 1)) / n
y = Inches(2.4)
x = Inches(0.5)
for i, (label, sub) in enumerate(steps):
    bx = x + (box_w + gap) * i
    color = TEAL if i % 2 == 0 else NAVY
    add_rect(s, bx, y, box_w, Inches(1.7), color)
    add_text(s, bx, y + Inches(0.2), box_w, Inches(0.45),
             f"{i+1:02d}", size=18, bold=True, color=SAND,
             align=PP_ALIGN.CENTER)
    add_text(s, bx, y + Inches(0.65), box_w, Inches(0.4),
             label, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, bx, y + Inches(1.05), box_w, Inches(0.6),
             sub, size=10, color=SAND, align=PP_ALIGN.CENTER)
    if i < n - 1:
        # arrow gap is too thin for an arrow shape; we just leave the gap
        pass

add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
         "What the user does vs. what the system does",
         size=14, bold=True, color=NAVY)
add_card(s, Inches(0.5), Inches(5.0), Inches(6.1), Inches(2.0),
         "USER", [
             "Defines brand profile and area constraints",
             "Adds existing branches (cannibalization)",
             "Picks target districts (optional)",
             "Reviews verdict + memo, compares finalists",
         ])
add_card(s, Inches(6.8), Inches(5.0), Inches(6.1), Inches(2.0),
         "SYSTEM", [
             "Builds candidate pool from listings + parcels",
             "Applies hard-fail and advisory gates",
             "Computes weighted scores + cannibalization + value",
             "Generates structured decision memo (LLM, cached)",
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 4 — Data sources overview
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "03 · DATA", "Data sources powering the Advisor", "4")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Two layers: persisted search/candidate state, and refreshed "
         "Riyadh-specific context tables.",
         size=15, color=INK)

# Two-column layout
add_card(s, Inches(0.5), Inches(2.6), Inches(6.1), Inches(4.4),
         "PERSISTED STATE  (per search)",
         [
             ("expansion_search", "brand, category, bounds, request JSON"),
             ("expansion_search_brand_profile", "price tier, sensitivity, goal"),
             ("expansion_search_existing_branch", "user branches for cannibalization"),
             ("expansion_candidate", "scores, gates, snapshots, memo, rerank"),
             ("expansion_candidate_feature_snapshot", "context cache per candidate"),
             ("expansion_saved_search", "title, filters, selected candidates, UI state"),
         ])
add_card(s, Inches(6.8), Inches(2.6), Inches(6.1), Inches(4.4),
         "CONTEXT  (refreshed weekly via GitHub Actions)",
         [
             ("expansion_road_context", "OSM roads → access & visibility"),
             ("expansion_parking_asset", "OSM parking → parking_score"),
             ("expansion_delivery_market", "Hungerstation/Jahez/Mrsool"),
             ("expansion_rent_comp", "rent comparables by district & asset"),
             ("expansion_competitor_quality", "POI + delivery + Google reviews"),
             ("riyadh_parcels_arcgis_proxy", "base parcel pool (Tier 3)"),
         ])
add_footer(s, "Sources fall back safely when context tables are empty; provenance kept in feature_snapshot.context_sources")

# ────────────────────────────────────────────────────────────────────────────
# Slide 5 — Context tables in detail
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "03 · DATA · CONTEXT", "Context tables — what each one feeds", "5")

rows = [
    ("expansion_road_context", "OSM (planet_osm_line / roads)",
     "road class, frontage_length_m, corner_lot, major-road touch",
     "access_score, frontage_score, visibility"),
    ("expansion_parking_asset", "OSM polygon / point",
     "amenity_type, capacity, walk_access_score, dropoff_score",
     "parking_score, parking_pass gate"),
    ("expansion_delivery_market", "delivery_source_record (normalized)",
     "platform, brand, category, rating, ETA, late-night flag",
     "delivery_demand, multi-platform presence"),
    ("expansion_rent_comp", "rent_comp + Aqar/Kaggle/CSV",
     "district, asset_type, rent_sar_m2_year",
     "rent estimation, rent_burden, value chip"),
    ("expansion_competitor_quality", "restaurant_poi + delivery + Google",
     "chain_strength, review_score, delivery_presence, overall",
     "competition_whitespace, comparable competitors"),
]

# header
top = Inches(2.05)
col_x = [Inches(0.5), Inches(3.3), Inches(5.7), Inches(8.9)]
col_w = [Inches(2.7), Inches(2.3), Inches(3.1), Inches(3.9)]
header_h = Inches(0.45)
add_rect(s, Inches(0.5), top, Inches(12.4), header_h, NAVY)
headers = ["Table", "Source", "Key fields", "Used by"]
for cx, cw, h in zip(col_x, col_w, headers):
    add_text(s, cx + Inches(0.1), top + Inches(0.07), cw - Inches(0.1),
             header_h, h, size=12, bold=True, color=SAND)

row_h = Inches(0.85)
for i, row in enumerate(rows):
    ry = top + header_h + row_h * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.5), ry, Inches(12.4), row_h, bg)
    for cx, cw, val, j in zip(col_x, col_w, row, range(4)):
        add_text(s, cx + Inches(0.1), ry + Inches(0.1), cw - Inches(0.1),
                 row_h - Inches(0.15), val,
                 size=11, bold=(j == 0),
                 color=NAVY if j == 0 else INK)

add_footer(s, "Each ingestion job runs on its own day; failures fall back to legacy sources, never silently degrade scoring.")

# ────────────────────────────────────────────────────────────────────────────
# Slide 6 — Ingestion schedule
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "03 · DATA · REFRESH", "Weekly ingestion schedule", "6")

add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Each context table is refreshed by a dedicated GitHub Actions "
         "workflow. All can also be triggered manually.",
         size=14, color=INK)

days = [
    ("MON 03:00 UTC", "Roads & Access",
     "expansion-advisor-data-roads.yml", "expansion_road_context"),
    ("TUE 04:00 UTC", "Parking",
     "expansion-advisor-data-parking.yml", "expansion_parking_asset"),
    ("WED 05:00 UTC", "Delivery Marketplace",
     "expansion-advisor-data-delivery.yml", "expansion_delivery_market"),
    ("THU 06:00 UTC", "Rent & Lease Comps",
     "expansion-advisor-data-rent-comps.yml", "expansion_rent_comp"),
    ("FRI 07:00 UTC", "Competitor Quality",
     "expansion-advisor-data-competitors.yml", "expansion_competitor_quality"),
]
y = Inches(2.7)
for i, (when, name, wf, tbl) in enumerate(days):
    ry = y + Inches(0.85) * i
    add_rect(s, Inches(0.5), ry, Inches(2.0), Inches(0.7), TEAL)
    add_text(s, Inches(0.5), ry + Inches(0.18), Inches(2.0), Inches(0.4),
             when, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(2.55), ry, Inches(10.4), Inches(0.7), LIGHT)
    add_text(s, Inches(2.7), ry + Inches(0.05), Inches(4.5), Inches(0.3),
             name, size=13, bold=True, color=NAVY)
    add_text(s, Inches(2.7), ry + Inches(0.36), Inches(5.0), Inches(0.3),
             wf, size=10, color=SUBTLE)
    add_text(s, Inches(7.8), ry + Inches(0.18), Inches(5.0), Inches(0.4),
             "→  " + tbl, size=12, bold=True, color=ACCENT)

add_footer(s, "Verify: SELECT COUNT(*) per table  ·  Provenance recorded in feature_snapshot.context_sources")

# ────────────────────────────────────────────────────────────────────────────
# Slide 7 — Candidate pool & enrichment
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "04 · PIPELINE", "Candidate pool and context enrichment", "7")

add_card(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(2.5),
         "TIERED CANDIDATE POOL",
         [
             ("Tier 1", "Aqar listings with rent (preferred)"),
             ("Tier 2", "occupied delivery POI"),
             ("Tier 3", "ArcGIS parcels (riyadh_parcels_arcgis_proxy)"),
             ("Stratification", "≤ 200 per district, max 2,000 total"),
             ("Headroom", "3× pool kept for scoring & rerank"),
         ])
add_card(s, Inches(6.8), Inches(2.0), Inches(6.1), Inches(2.5),
         "ENRICHMENT (LATERAL JOINs)",
         [
             ("Roads", "distance to major road, frontage, corner"),
             ("Parking", "nearby capacity, walk-access score"),
             ("Population", "H3 hex catchment density"),
             ("Delivery", "listing count, rating, ETA, platforms"),
             ("Competitors", "nearby chain branches, strength"),
         ])
add_card(s, Inches(0.5), Inches(4.65), Inches(12.4), Inches(2.4),
         "PROVENANCE  (feature_snapshot.context_sources)",
         [
             "road_source: 'expansion_road_context' or 'estimated'",
             "parking_source: 'expansion_parking_asset' or 'estimated'",
             "delivery_source: 'expansion_delivery_market' or 'delivery_source_record'",
             "rent_source: 'expansion_rent_district_retail' / 'aqar_district' / city-level fallback",
             "competitor_source: 'expansion_competitor_quality' or 'restaurant_poi'",
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 8 — Gates
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "05 · GATES", "Pass / fail filters before scoring", "8")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Gates are evaluated per candidate. Hard-fail gates exclude a "
         "candidate; advisory gates only annotate.",
         size=14, color=INK)

add_card(s, Inches(0.5), Inches(2.7), Inches(4.0), Inches(4.5),
         "STRUCTURAL  (always hard-fail)",
         [
             ("zoning_fit_pass", "landuse acceptable for F&B"),
             ("area_fit_pass", "unit area within min/max m²"),
         ])
add_card(s, Inches(4.7), Inches(2.7), Inches(4.0), Inches(4.5),
         "OPERATIONAL  (configurable)",
         [
             ("frontage_access_pass", "≥ 6m or near major road"),
             ("parking_pass", "min capacity / walk score"),
             ("district_pass", "matches target districts"),
             ("cannibalization_pass", "> tolerance from branches"),
             ("delivery_market_pass", "delivery listings present"),
             ("economics_pass", "rent burden ≥ 20"),
         ])
add_card(s, Inches(8.9), Inches(2.7), Inches(4.0), Inches(4.5),
         "VIABILITY FLOORS  &  ADVISORY",
         [
             ("population_floor_pass", "min residents in catchment"),
             ("commercial_floor_pass", "min brand presence"),
             ("construction_proximity_pass", "buffer from construction"),
             ("radiance_growth_pass", "advisory only (NASA Black Marble)"),
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 9 — Scoring weights
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "06 · SCORING", "Six weighted dimensions → final score (0–100)", "9")

dims = [
    ("Demand potential", 25, "population, foot traffic, district momentum", TEAL),
    ("Competition whitespace", 20, "inverse competitor density", NAVY),
    ("Occupancy economics", 20, "rent burden vs. estimated revenue", ACCENT),
    ("Delivery demand", 15, "marketplace size, category, rating", TEAL),
    ("Access & visibility", 10, "frontage, road, parking", NAVY),
    ("Brand fit", 10, "districts, expansion goal, cannibalization", ACCENT),
]
y = Inches(2.2)
for i, (name, w, desc, color) in enumerate(dims):
    ry = y + Inches(0.7) * i
    # name
    add_text(s, Inches(0.5), ry, Inches(3.0), Inches(0.55),
             name, size=14, bold=True, color=NAVY)
    # bar
    bar_x = Inches(3.6)
    bar_max_w = Inches(5.5)
    bar_w = Inches(5.5 * (w / 25))
    add_rect(s, bar_x, ry + Inches(0.12), bar_max_w, Inches(0.32), LIGHT)
    add_rect(s, bar_x, ry + Inches(0.12), bar_w, Inches(0.32), color)
    add_text(s, bar_x + bar_w + Inches(0.1), ry + Inches(0.08),
             Inches(0.9), Inches(0.4),
             f"{w}%", size=14, bold=True, color=color)
    # desc
    add_text(s, Inches(10.0), ry + Inches(0.1), Inches(3.0),
             Inches(0.5), desc, size=11, color=SUBTLE)

add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
         "Cannibalization and value are computed alongside but not part of "
         "the weighted sum — they refine ranking and surface chips in the UI.",
         size=11, color=SUBTLE)

# ────────────────────────────────────────────────────────────────────────────
# Slide 10 — Service models & saturation
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "06 · SCORING · DEMAND",
             "Population score depends on the service model", "10")

models = [
    ("QSR", "1.5 km", "80,000", TEAL),
    ("Café", "1.0 km", "40,000", ACCENT),
    ("Dine-in", "3.5 km", "250,000", NAVY),
    ("Delivery-first", "3.0 km", "180,000", TEAL),
]
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "The catchment radius and saturation reference change by format. "
         "Score saturates at the residents listed below.",
         size=14, color=INK)

card_w = Inches(3.0)
gap = Inches(0.13)
y = Inches(2.9)
for i, (name, radius, pop, color) in enumerate(models):
    x = Inches(0.5) + (card_w + gap) * i
    add_rect(s, x, y, card_w, Inches(2.6), color)
    add_text(s, x, y + Inches(0.2), card_w, Inches(0.5),
             name, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.95), card_w, Inches(0.4),
             "Catchment", size=11, color=SAND, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.2), card_w, Inches(0.5),
             radius, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.75), card_w, Inches(0.3),
             "Saturation reference", size=11, color=SAND,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.05), card_w, Inches(0.45),
             pop, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_footer(s, "Population reach pulls from H3 hex density aggregations across the catchment radius.")

# ────────────────────────────────────────────────────────────────────────────
# Slide 11 — Cannibalization, economics & value
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "06 · SCORING · DERIVED",
             "Cannibalization, economics, and the value chip", "11")

add_card(s, Inches(0.5), Inches(2.0), Inches(4.0), Inches(5.0),
         "CANNIBALIZATION",
         [
             ("Distance-based", "to nearest user-supplied branch"),
             ("Tolerance", "default 1,800m (range 800–2,500m)"),
             ("Feeds brand_fit", "and the cannibalization gate"),
             ("Goal aware", "preferred spacing varies by expansion goal"),
         ])
add_card(s, Inches(4.7), Inches(2.0), Inches(4.0), Inches(5.0),
         "ECONOMICS",
         [
             ("rent_burden_score", "100 × est. revenue / annual rent"),
             ("Revenue model", "demand × category × location"),
             ("Three modes", "peer-relative · absolute envelope · legacy"),
             ("economics_pass", "fails when rent_burden < 20"),
         ])
add_card(s, Inches(8.9), Inches(2.0), Inches(4.0), Inches(5.0),
         "VALUE CHIP",
         [
             ("value_score", "geometric mean of revenue & rent burden"),
             ("best_value", "top decile chip"),
             ("neutral", "middle 80%"),
             ("above_market", "bottom decile (downrank candidate)"),
             ("Null", "when rent_burden ran in fallback mode"),
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 12 — Confidence grade
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "07 · CONFIDENCE",
             "How much can you trust each candidate?", "12")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Confidence grade reflects data completeness — independent of the "
         "score. Sparse signals don't lower the score, they lower trust.",
         size=14, color=INK)

grades = [
    ("A", "≥ 80%", "all major snapshots populated", GREEN),
    ("B", "60 – 80%", "most signals present", TEAL),
    ("C", "40 – 60%", "partial; some estimated values", ACCENT),
    ("D", "< 40%", "sparse / heavily estimated", RED),
]
y = Inches(2.9)
card_w = Inches(3.0)
gap = Inches(0.13)
for i, (g, pct, desc, color) in enumerate(grades):
    x = Inches(0.5) + (card_w + gap) * i
    add_rect(s, x, y, card_w, Inches(2.8), LIGHT)
    add_rect(s, x, y, card_w, Inches(0.85), color)
    add_text(s, x, y + Inches(0.18), card_w, Inches(0.6),
             f"GRADE  {g}", size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.05), card_w, Inches(0.5),
             "Completeness", size=11, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.3), card_w, Inches(0.5),
             pct, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(2.0), card_w - Inches(0.4),
             Inches(0.7), desc, size=12, color=INK, align=PP_ALIGN.CENTER)

add_footer(s, "Source: feature_snapshot.data_completeness_score across road, parking, delivery, population, competitor signals.")

# ────────────────────────────────────────────────────────────────────────────
# Slide 13 — LLM rerank & memo
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "08 · LLM LAYER",
             "Reranking and decision memos", "13")

add_card(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(5.0),
         "LLM RERANK  (Phase 2, optional)",
         [
             ("Toggle", "EXPANSION_LLM_RERANK_ENABLED"),
             ("Bounded shortlist", "top N candidates only"),
             ("Structured output", "rerank_reason, rerank_delta"),
             ("Fuzzy tiebreak", "within 1.5-pt bands prefer richer LLM signal"),
             ("Value-band reorder", "uprank best_value, downrank above_market"),
             ("Deterministic fallback", "always available (rerank_status='skipped')"),
         ])
add_card(s, Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.0),
         "DECISION MEMO  (LLM, structured)",
         [
             ("Pre-warmed", "background task generates top-N at search time"),
             ("Cached", "per search_id + parcel_id under MEMO_PROMPT_VERSION"),
             ("Structured JSON", "headline · evidence · risks · property · finance · market · competitive · bottom line"),
             ("Legacy fallback", "rendered text if structured path fails"),
             ("Daily ceiling", "LLM budget enforced (503 on breach)"),
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 14 — Dedup / shortlist diversity
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "09 · SHORTLIST",
             "Deduplication & shortlist diversity", "14")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Aggressive dedupe risks losing real options. The Advisor uses "
         "loose spatial + attribute keys to remove clones without trimming "
         "genuine variety.",
         size=14, color=INK)

add_card(s, Inches(0.5), Inches(2.7), Inches(6.1), Inches(4.3),
         "DEDUPE KEYS",
         [
             ("Spatial", "55m grid snap"),
             ("District", "normalized district key"),
             ("Area bucket", "size class"),
             ("Rent bucket", "price band"),
             ("Aggressive mode", "report view adds economics-similarity keys"),
         ])
add_card(s, Inches(6.8), Inches(2.7), Inches(6.1), Inches(4.3),
         "POOL HEADROOM",
         [
             ("Max from SQL", "2,000 candidates"),
             ("Stratification", "≤ 200 per district (geographic spread)"),
             ("3× headroom", "kept beyond final list for safe rerank"),
             ("Soft demote", "pillar legs reduce weak cohorts gracefully"),
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 15 — UI: Results panel
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "10 · UI", "Results panel — what the user sees first", "15")

# left panel mock
mock_x = Inches(0.5)
mock_y = Inches(2.0)
mock_w = Inches(6.5)
mock_h = Inches(5.0)
add_rect(s, mock_x, mock_y, mock_w, mock_h, LIGHT)
add_rect(s, mock_x, mock_y, mock_w, Inches(0.5), NAVY)
add_text(s, mock_x + Inches(0.15), mock_y + Inches(0.1),
         mock_w - Inches(0.3), Inches(0.3),
         "Expansion results · 8 pass · 12 candidates",
         size=12, bold=True, color=SAND)
# rows
for i in range(4):
    ry = mock_y + Inches(0.65) + Inches(1.05) * i
    add_rect(s, mock_x + Inches(0.15), ry, mock_w - Inches(0.3),
             Inches(0.95), WHITE)
    # rank badge
    add_rect(s, mock_x + Inches(0.3), ry + Inches(0.15),
             Inches(0.5), Inches(0.5), TEAL)
    add_text(s, mock_x + Inches(0.3), ry + Inches(0.25),
             Inches(0.5), Inches(0.4),
             f"{i+1}", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # title + district
    add_text(s, mock_x + Inches(0.95), ry + Inches(0.1),
             Inches(3.5), Inches(0.3),
             f"Candidate {i+1}  ·  Al Olaya",
             size=12, bold=True, color=NAVY)
    add_text(s, mock_x + Inches(0.95), ry + Inches(0.4),
             Inches(4.5), Inches(0.3),
             "320 m²  ·  430k SAR/yr  ·  confidence A",
             size=10, color=SUBTLE)
    # score pill
    add_rect(s, mock_x + Inches(5.4), ry + Inches(0.2),
             Inches(0.85), Inches(0.4), NAVY)
    add_text(s, mock_x + Inches(5.4), ry + Inches(0.27),
             Inches(0.85), Inches(0.3),
             f"{82 - i*4}.{i}", size=12, bold=True, color=SAND,
             align=PP_ALIGN.CENTER)
    # gate badge
    color = GREEN if i < 3 else GREY
    label = "PASS" if i < 3 else "PARTIAL"
    add_rect(s, mock_x + Inches(0.95), ry + Inches(0.65),
             Inches(0.7), Inches(0.22), color)
    add_text(s, mock_x + Inches(0.95), ry + Inches(0.66),
             Inches(0.7), Inches(0.22),
             label, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # value chip
    color = ACCENT if i == 0 else (TEAL if i < 3 else GREY)
    label = "best_value" if i == 0 else ("neutral" if i < 3 else "above_market")
    add_rect(s, mock_x + Inches(1.7), ry + Inches(0.65),
             Inches(1.3), Inches(0.22), color)
    add_text(s, mock_x + Inches(1.7), ry + Inches(0.66),
             Inches(1.3), Inches(0.22),
             label, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# right side: legend
add_card(s, Inches(7.3), Inches(2.0), Inches(5.6), Inches(5.0),
         "WHAT EACH ROW TELLS YOU",
         [
             ("Rank badge", "1-based final rank (post-rerank)"),
             ("Score pill", "weighted composite (0–100)"),
             ("Gate badge", "PASS / FAIL / PARTIAL"),
             ("Value chip", "best_value · neutral · above_market"),
             ("Confidence", "A–D letter grade"),
             ("Quick facts", "area · annual rent · district"),
             ("Click", "opens the Decision Memory drawer"),
         ])
add_footer(s, "Summary strip on top: pass count · best overall · best value highlight  ·  filters: score/gate/district/area")

# ────────────────────────────────────────────────────────────────────────────
# Slide 16 — Decision Memory drawer (Memo tab)
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "11 · UI", "Decision Memory drawer — Memo tab", "16")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Default tab. Narrative for stakeholders — readable end-to-end "
         "without diving into raw numbers.",
         size=14, color=INK)

# verdict row mock
add_rect(s, Inches(0.5), Inches(2.6), Inches(12.4), Inches(0.7), NAVY)
add_text(s, Inches(0.7), Inches(2.7), Inches(2.5), Inches(0.5),
         "VERDICT", size=11, color=SAND)
add_rect(s, Inches(0.7), Inches(2.95), Inches(1.2), Inches(0.3), GREEN)
add_text(s, Inches(0.7), Inches(2.96), Inches(1.2), Inches(0.3),
         "GO", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.4), Inches(2.7), Inches(1.5), Inches(0.5),
         "CONFIDENCE  A", size=12, bold=True, color=SAND)
add_text(s, Inches(4.5), Inches(2.7), Inches(2.5), Inches(0.5),
         "SCORE  82.4", size=12, bold=True, color=SAND)
add_text(s, Inches(7.5), Inches(2.7), Inches(5.0), Inches(0.5),
         "AL OLAYA  ·  320 m²  ·  430k SAR/yr",
         size=12, color=SAND, align=PP_ALIGN.RIGHT)

# memo sections
sections = [
    ("Headline recommendation", "Strong demand at fair price."),
    ("Key evidence", "polarity-tagged positives / neutrals / negatives"),
    ("Property overview", "area · frontage · street type · parking · listing age"),
    ("Financial framing", "annual rent · comp median · percentile · spread"),
    ("Market context", "population reach · district momentum · realized 30d demand"),
    ("Competitive landscape", "top chains · distances · saturation thesis"),
    ("Risks", "structured list with mitigations"),
    ("Bottom line", "recommended next action"),
]
y = Inches(3.6)
col1_x = Inches(0.5)
col2_x = Inches(6.85)
col_w = Inches(6.05)
for i, (h, d) in enumerate(sections):
    x = col1_x if i < 4 else col2_x
    ry = y + Inches(0.7) * (i % 4)
    add_rect(s, x, ry, col_w, Inches(0.6), LIGHT)
    add_rect(s, x, ry, Inches(0.08), Inches(0.6), TEAL)
    add_text(s, x + Inches(0.2), ry + Inches(0.05), col_w - Inches(0.3),
             Inches(0.3), h, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), ry + Inches(0.3), col_w - Inches(0.3),
             Inches(0.3), d, size=10, color=SUBTLE)

add_footer(s, "Lead candidates also expose a Copy-Ready Summary block for stakeholder sharing.")

# ────────────────────────────────────────────────────────────────────────────
# Slide 17 — Diagnostics tab
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "11 · UI", "Decision Memory drawer — Diagnostics tab", "17")
add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
         "Five inner sub-tabs and a full Decision Logic card. This is where "
         "an analyst checks why a candidate scored what it did.",
         size=14, color=INK)

subtabs = [
    ("Breakdown", "stacked weighted components + raw JSON"),
    ("Site Fit", "road, parking, visibility, observed vs. estimated"),
    ("Delivery", "platform mix, category penetration, rating, ETA"),
    ("Competition", "nearby chains, distances, saturation tier"),
    ("Economics", "rent comps, percentile, peer range, revenue model"),
]
y = Inches(2.7)
sw = Inches(2.42)
for i, (name, desc) in enumerate(subtabs):
    x = Inches(0.5) + (sw + Inches(0.1)) * i
    add_rect(s, x, y, sw, Inches(0.55), TEAL if i == 0 else NAVY)
    add_text(s, x, y + Inches(0.1), sw, Inches(0.4),
             name, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.65), sw, Inches(1.0),
             desc, size=10, color=SUBTLE, align=PP_ALIGN.CENTER)

# decision logic card
add_card(s, Inches(0.5), Inches(4.6), Inches(6.1), Inches(2.5),
         "FULL SCORE BREAKDOWN  (collapsible)",
         [
             ("weights", "component weights"),
             ("inputs", "raw feature values"),
             ("weighted_components", "computed scores per dimension"),
             ("display", "human-friendly labels"),
             ("economics_detail", "rent burden, value_score, value_band"),
         ])
add_card(s, Inches(6.8), Inches(4.6), Inches(6.1), Inches(2.5),
         "DECISION LOGIC CARD",
         [
             ("Gate status", "color-coded pass / fail / unknown"),
             ("Per gate", "name, threshold, actual, explanation"),
             ("Pool diagnostics", "hard-floor drops, demote-leg triggers"),
             ("Hard vs. advisory", "advisory gates never block overall_pass"),
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 18 — Compare & Report
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "11 · UI", "Compare and Report views", "18")

add_card(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(5.0),
         "COMPARE PANEL  (2–6 candidates)",
         [
             ("Side-by-side", "score, gates, economics, demand, delivery"),
             ("best_overall", "highest final_score"),
             ("lowest_cannibalization", "furthest from existing branches"),
             ("highest_demand", "best demand_potential"),
             ("best_economics", "lowest rent_burden / best peer percentile"),
             ("strongest_whitespace", "least competitive saturation"),
             ("best_value", "best revenue-to-rent geometric mean"),
             ("most_confident", "highest data completeness"),
         ])
add_card(s, Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.0),
         "RECOMMENDATION REPORT",
         [
             ("Top 3", "featured with quick facts and breakdowns"),
             ("Recommendation summary",
              "best_candidate_id, runner_up, why_best, main_risk, best_format"),
             ("pass_count", "candidates clearing all hard-fail gates"),
             ("Dimension winners", "for each scored dimension"),
             ("Assumptions", "city: Riyadh · parcel_source: listings_only"),
         ])

# ────────────────────────────────────────────────────────────────────────────
# Slide 19 — Glossary 1
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "12 · GLOSSARY", "Plain-English definitions (1 of 2)", "19")

g1 = [
    ("Cannibalization",
     "Market overlap between a new location and your existing branches; "
     "controlled by tolerance distance and the cannibalization gate."),
    ("Confidence grade",
     "A–D rating of data quality. A = rich signals, D = sparse / "
     "estimated. Independent of the score."),
    ("Data completeness score",
     "Integer 0–100 measuring availability of road, parking, delivery, "
     "and population signals. Drives the confidence grade."),
    ("Decision memo",
     "Structured LLM narrative explaining a candidate. Cached per "
     "search and parcel."),
    ("Delivery market",
     "Aggregated presence on Hungerstation, Jahez, Mrsool — multi-platform "
     "saturation and category penetration."),
    ("Economics score",
     "Rental viability metric: estimated annual revenue vs. annual rent. "
     "Drives economics_pass and the value chip."),
    ("Feature snapshot",
     "Lightweight cache of road, parking, population, delivery, and "
     "competitor signals; tracks observed vs. estimated."),
    ("Gate",
     "Pass/fail filter applied before scoring. Structural gates are "
     "always hard-fail; operational gates can be advisory."),
]
y = Inches(1.95)
row_h = Inches(0.62)
for i, (term, defn) in enumerate(g1):
    ry = y + row_h * i
    add_rect(s, Inches(0.5), ry + Inches(0.05), Inches(2.7),
             row_h - Inches(0.1), NAVY)
    add_text(s, Inches(0.6), ry + Inches(0.13), Inches(2.6),
             row_h - Inches(0.2), term, size=12, bold=True, color=SAND)
    add_text(s, Inches(3.4), ry + Inches(0.08), Inches(9.5),
             row_h - Inches(0.1), defn, size=12, color=INK)

# ────────────────────────────────────────────────────────────────────────────
# Slide 20 — Glossary 2
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "12 · GLOSSARY", "Plain-English definitions (2 of 2)", "20")

g2 = [
    ("Rank position",
     "1-based position in the final list. deterministic_rank is pre-LLM, "
     "final_rank is post-LLM."),
    ("Score breakdown",
     "Per-component decomposition (weights, inputs, weighted_components, "
     "labels, final_score). Surfaced in Diagnostics."),
    ("Service model",
     "Operating format: QSR, dine-in, delivery-first, café. Drives "
     "catchment radius and saturation thresholds."),
    ("Shortlist diversity",
     "Dedupe logic that removes spatial/attribute clones without "
     "trimming genuine variety."),
    ("Value band",
     "best_value · neutral · above_market chip from the geometric mean "
     "of revenue index and rent burden."),
    ("Verdict",
     "GO / CAUTION / RECONSIDER summary written into the memo."),
    ("Whitespace",
     "Inverse competitor density. competition_whitespace = 100 means "
     "no competitors; 0 means saturated."),
    ("LLM rerank",
     "Optional structured reranking of the top shortlist by Claude, "
     "with rerank_reason and rerank_delta."),
]
y = Inches(1.95)
row_h = Inches(0.62)
for i, (term, defn) in enumerate(g2):
    ry = y + row_h * i
    add_rect(s, Inches(0.5), ry + Inches(0.05), Inches(2.7),
             row_h - Inches(0.1), NAVY)
    add_text(s, Inches(0.6), ry + Inches(0.13), Inches(2.6),
             row_h - Inches(0.2), term, size=12, bold=True, color=SAND)
    add_text(s, Inches(3.4), ry + Inches(0.08), Inches(9.5),
             row_h - Inches(0.1), defn, size=12, color=INK)

# ────────────────────────────────────────────────────────────────────────────
# Slide 21 — API endpoints
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "13 · API", "Backend endpoints under /v1/expansion-advisor", "21")

groups = [
    ("Search lifecycle", [
        "POST   /searches",
        "GET    /searches/{id}",
        "GET    /searches/{id}/candidates",
        "GET    /searches/{id}/report",
    ]),
    ("Candidate detail", [
        "GET    /candidates/{id}/memo",
        "POST   /candidates/compare",
        "POST   /decision-memo",
    ]),
    ("Saved searches", [
        "POST   /saved-searches",
        "GET    /saved-searches",
        "GET    /saved-searches/{id}",
        "PATCH  /saved-searches/{id}",
        "DELETE /saved-searches/{id}",
    ]),
    ("Supporting", [
        "GET    /districts",
        "GET    /branch-suggestions",
    ]),
]
y = Inches(2.1)
card_w = Inches(6.1)
card_h = Inches(2.4)
positions = [
    (Inches(0.5), y),
    (Inches(6.8), y),
    (Inches(0.5), y + card_h + Inches(0.2)),
    (Inches(6.8), y + card_h + Inches(0.2)),
]
for (title, lines), (cx, cy) in zip(groups, positions):
    add_rect(s, cx, cy, card_w, card_h, LIGHT)
    add_rect(s, cx, cy, card_w, Inches(0.4), NAVY)
    add_text(s, cx + Inches(0.15), cy + Inches(0.06),
             card_w - Inches(0.3), Inches(0.3),
             title, size=12, bold=True, color=SAND)
    for j, line in enumerate(lines):
        add_text(s, cx + Inches(0.2), cy + Inches(0.5) + Inches(0.32) * j,
                 card_w - Inches(0.3), Inches(0.3),
                 line, size=12, color=INK, font="Consolas")

add_footer(s, "All list responses use { items: [...] }  ·  Memo response includes structured + legacy fields")

# ────────────────────────────────────────────────────────────────────────────
# Slide 22 — Recap / Where to look
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "14 · RECAP", "Putting it all together", "22")

add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
         "If you remember five things…", size=18, bold=True, color=NAVY)
recap = [
    ("Riyadh-only, listings-first.",
     "Pool comes from Aqar / Wasalt / Bayut listings, with parcel + delivery "
     "POI fallbacks; everything is stratified by district."),
    ("Gates first, scoring second.",
     "Hard-fail gates prune the pool. Six weighted dimensions then produce "
     "a 0–100 score. Cannibalization, economics, and value refine it."),
    ("Two layers in the drawer.",
     "Memo for stakeholders (verdict + narrative). Diagnostics for analysts "
     "(weights, inputs, gates, comps)."),
    ("Trust ≠ score.",
     "Confidence grade A–D tells you how complete the data is, separately "
     "from the score itself."),
    ("LLM is bounded.",
     "Reranking and memos are structured, cached, and budget-capped. "
     "Deterministic results are always available as a fallback."),
]
y = Inches(2.6)
for i, (head, body) in enumerate(recap):
    ry = y + Inches(0.85) * i
    add_rect(s, Inches(0.5), ry + Inches(0.1), Inches(0.5),
             Inches(0.5), TEAL)
    add_text(s, Inches(0.5), ry + Inches(0.18), Inches(0.5),
             Inches(0.5), f"{i+1}", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.15), ry + Inches(0.08), Inches(11.5),
             Inches(0.35), head, size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.15), ry + Inches(0.4), Inches(11.5),
             Inches(0.5), body, size=12, color=INK)

# ────────────────────────────────────────────────────────────────────────────
# Slide 23 — Where things live (file map)
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_header(s, "15 · CODE MAP", "Where to look in the repo", "23")

paths = [
    ("Backend service", "app/services/expansion_advisor.py"),
    ("Backend API", "app/api/expansion_advisor.py"),
    ("Ingestion jobs", "app/ingest/expansion_advisor_*.py"),
    ("Migrations", "alembic/versions/20260310_*  →  20260314_*"),
    ("Frontend feature", "frontend/src/features/expansion-advisor/"),
    ("Memo panel", "frontend/src/features/.../ExpansionMemoPanel.tsx"),
    ("Results panel", "frontend/src/features/.../ExpansionResultsPanel.tsx"),
    ("Workflows", ".github/workflows/expansion-advisor-data-*.yml"),
    ("Docs", "docs/expansion_advisor_data_ingest.md"),
]
y = Inches(2.1)
row_h = Inches(0.5)
for i, (label, path) in enumerate(paths):
    ry = y + row_h * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.5), ry, Inches(12.4), row_h, bg)
    add_text(s, Inches(0.7), ry + Inches(0.1), Inches(3.6),
             Inches(0.4), label, size=13, bold=True, color=NAVY)
    add_text(s, Inches(4.4), ry + Inches(0.1), Inches(8.4),
             Inches(0.4), path, size=12, color=INK, font="Consolas")

add_footer(s, "Service version: expansion_advisor_v7  ·  Riyadh only  ·  EPSG:4326 (metric via EPSG:32638)")

# ────────────────────────────────────────────────────────────────────────────
# Slide 24 — Closing
# ────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, Inches(3.4), SLIDE_W, Inches(0.06), TEAL)
add_text(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(1.0),
         "Questions?", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.6), Inches(11.9), Inches(0.6),
         "The Expansion Advisor — data, gates, scoring, memo, and UI.",
         size=20, color=SAND)
add_text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.4),
         "For deeper drills:  /docs/expansion_advisor_data_ingest.md",
         size=14, color=SAND)
add_text(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(0.4),
         "Service:  app/services/expansion_advisor.py",
         size=14, color=SAND)
add_text(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.4),
         "Frontend:  frontend/src/features/expansion-advisor/",
         size=14, color=SAND)
add_text(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.4),
         "OAKTREE ATLAS  ·  RIYADH",
         size=12, bold=True, color=ACCENT)

prs.save(OUT)
print(f"Wrote {OUT}  ·  {len(prs.slides)} slides")
