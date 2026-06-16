"""
Generates a non-technical PowerPoint user guide for the Expansion Advisor.

Audience: F&B operators (restaurant/cafe owners and chain managers) with no
technical background who want to find new branch locations in Riyadh.

Run:
    python docs/generate_expansion_advisor_guide.py

Output:
    docs/Expansion-Advisor-User-Guide.pptx
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "assets", "expansion-advisor")

# ---------------------------------------------------------------------------
# Brand palette (calm, professional, F&B-friendly)
# ---------------------------------------------------------------------------
OAK_GREEN = RGBColor(0x1F, 0x4D, 0x3E)   # deep oak green
LEAF = RGBColor(0x3C, 0x8C, 0x5A)        # accent green
SAND = RGBColor(0xF4, 0xEF, 0xE6)        # warm light background
INK = RGBColor(0x23, 0x2A, 0x2E)         # near-black text
SLATE = RGBColor(0x5C, 0x67, 0x6B)       # muted secondary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xC8, 0x8A, 0x1E)       # caution / "unknown"
RED = RGBColor(0xB3, 0x3A, 0x3A)         # fail
GREEN_OK = RGBColor(0x2E, 0x8B, 0x57)    # pass
CARD = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def add_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    bg_shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W, SLIDE_H  # 1 = rectangle
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg_shape._element)
    spTree.insert(2, bg_shape._element)
    return slide


def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1), radius=False):
    shape_type = 5 if radius else 1  # 5 = rounded rectangle
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, space_after=Pt(6), line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = size
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = "Calibri"
    return tb


def P(txt, size=18, color=INK, bold=False, italic=False):
    """Shorthand to build a single-run paragraph."""
    return [(txt, Pt(size), color, bold, italic)]


def section_header(slide, kicker, title):
    """Standard content-slide header band."""
    rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=OAK_GREEN)
    rect(slide, Inches(0.55), Inches(0.32), Inches(0.12), Inches(0.62), fill=LEAF)
    text(slide, Inches(0.85), Inches(0.18), Inches(11.8), Inches(0.4),
         [P(kicker, 13, RGBColor(0xBF, 0xD9, 0xC8), bold=True)],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, Inches(0.85), Inches(0.5), Inches(11.8), Inches(0.6),
         [P(title, 28, WHITE, bold=True)], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n):
    text(slide, Inches(0.55), Inches(7.02), Inches(8), Inches(0.35),
         [P("Oaktree Atlas  ·  Expansion Advisor User Guide", 10, SLATE)])
    text(slide, Inches(11.0), Inches(7.02), Inches(1.8), Inches(0.35),
         [P(str(n), 10, SLATE, bold=True)], align=PP_ALIGN.RIGHT)


def bullet_block(slide, x, y, w, items, size=16, gap=Pt(10), color=INK,
                 marker="●", marker_color=LEAF):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            paras.append([(f"{marker}  ", Pt(size), marker_color, True, False),
                          (head, Pt(size), color, True, False),
                          (f"  {body}", Pt(size), SLATE, False, False)])
        else:
            paras.append([(f"{marker}  ", Pt(size), marker_color, True, False),
                          (it, Pt(size), color, False, False)])
    return text(slide, x, y, w, Inches(5), paras, space_after=gap, line_spacing=1.05)


def card(slide, x, y, w, h, title, body_lines, accent=LEAF, title_size=16,
         body_size=13):
    rect(slide, x, y, w, h, fill=CARD, line=RGBColor(0xE2, 0xDD, 0xD2), line_w=Pt(1),
         radius=True)
    rect(slide, x, y, Inches(0.10), h, fill=accent, radius=True)
    text(slide, x + Inches(0.28), y + Inches(0.14), w - Inches(0.45), Inches(0.5),
         [P(title, title_size, OAK_GREEN, bold=True)])
    paras = [P(line, body_size, SLATE) for line in body_lines]
    text(slide, x + Inches(0.28), y + Inches(0.62), w - Inches(0.45), h - Inches(0.7),
         paras, space_after=Pt(4), line_spacing=1.0)


def screenshot_slide(kicker, title, img_file, blurb, annotations,
                     frame_max_h=5.35, frame_max_w=5.6):
    """A slide pairing a real product screenshot with plain-English callouts."""
    s = add_slide(SAND)
    section_header(s, kicker, title)

    path = os.path.join(SHOTS, img_file)
    iw, ih = Image.open(path).size
    aspect = iw / ih
    # Fit the screenshot within the frame box, preserving aspect ratio.
    h = frame_max_h
    w = h * aspect
    if w > frame_max_w:
        w = frame_max_w
        h = w / aspect
    img_x = Inches(0.7)
    img_y = Inches(1.45) + Inches((frame_max_h - h) / 2)
    # Soft "screen" frame behind the screenshot.
    rect(s, img_x - Inches(0.12), Inches(1.4), Inches(w) + Inches(0.24),
         Inches(frame_max_h + 0.1), fill=WHITE, line=RGBColor(0xE2, 0xDD, 0xD2),
         line_w=Pt(1.2), radius=True)
    s.shapes.add_picture(path, img_x, img_y, height=Inches(h))

    # Annotation column to the right of the screenshot.
    ax = img_x + Inches(w) + Inches(0.55)
    aw = SLIDE_W - ax - Inches(0.55)
    text(s, ax, Inches(1.5), aw, Inches(0.9), [P(blurb, 15, INK)], line_spacing=1.1)
    ay = Inches(2.5)
    for head, body, accent in annotations:
        rect(s, ax, ay, aw, Inches(0.86), fill=WHITE,
             line=RGBColor(0xE2, 0xDD, 0xD2), radius=True)
        rect(s, ax, ay, Inches(0.1), Inches(0.86), fill=accent, radius=True)
        text(s, ax + Inches(0.28), ay + Inches(0.05), aw - Inches(0.45), Inches(0.78),
             [[(head + "  ", Pt(15), OAK_GREEN, True, False),
               (body, Pt(13), SLATE, False, False)]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        ay += Inches(0.98)
    text(s, ax, Inches(6.5), aw, Inches(0.4),
         [P("Real Expansion Advisor screen · sample Riyadh data", 11, SLATE, italic=True)])
    footer(s, page())
    return s


PAGE = 0


def page():
    global PAGE
    PAGE += 1
    return PAGE


# ===========================================================================
# SLIDE 1 — Title
# ===========================================================================
s = add_slide(OAK_GREEN)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=OAK_GREEN)
# decorative side band
rect(s, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=LEAF)
rect(s, Inches(9.4), Inches(0), Inches(3.93), SLIDE_H, fill=RGBColor(0x1A, 0x42, 0x35))

text(s, Inches(0.9), Inches(1.7), Inches(8.2), Inches(0.5),
     [P("OAKTREE ATLAS", 18, RGBColor(0x9C, 0xC4, 0xAD), bold=True)])
text(s, Inches(0.9), Inches(2.25), Inches(8.4), Inches(2.2),
     [P("Expansion Advisor", 54, WHITE, bold=True),
      P("A simple guide to finding your next branch", 26, RGBColor(0xCF, 0xE3, 0xD7))],
     space_after=Pt(10))
text(s, Inches(0.9), Inches(4.7), Inches(8.0), Inches(1.0),
     [P("For restaurant & café operators in Riyadh", 18, RGBColor(0xBF, 0xD9, 0xC8))])

# right-band callouts
text(s, Inches(9.75), Inches(2.0), Inches(3.3), Inches(4),
     [P("Inside this guide", 15, WHITE, bold=True),
      P("• What it does", 14, RGBColor(0xCF, 0xE3, 0xD7)),
      P("• How to run a search", 14, RGBColor(0xCF, 0xE3, 0xD7)),
      P("• Reading your results", 14, RGBColor(0xCF, 0xE3, 0xD7)),
      P("• Comparing sites", 14, RGBColor(0xCF, 0xE3, 0xD7)),
      P("• Plain-English glossary", 14, RGBColor(0xCF, 0xE3, 0xD7))],
     space_after=Pt(10))
text(s, Inches(0.9), Inches(6.7), Inches(8), Inches(0.4),
     [P("No technical knowledge required.", 14, RGBColor(0x9C, 0xC4, 0xAD), italic=True)])

# ===========================================================================
# SLIDE 2 — What is the Expansion Advisor?
# ===========================================================================
s = add_slide(SAND)
section_header(s, "GETTING STARTED", "What is the Expansion Advisor?")
text(s, Inches(0.85), Inches(1.55), Inches(11.6), Inches(1.0),
     [P("It is your location scout. You tell it about your brand and where your "
        "current branches are. It then looks across hundreds of commercial sites "
        "in Riyadh and tells you which ones are the best places to open next — "
        "and why.", 18, INK)], line_spacing=1.15)

cards = [
    ("Finds sites for you", ["Scans the whole city so you", "don't have to drive around", "scouting blindly."], LEAF),
    ("Ranks them 1–15", ["Sorts the best candidates", "to the top with a clear", "score out of 100."], OAK_GREEN),
    ("Explains every result", ["Shows the reasons a site is", "strong or risky in plain", "language."], AMBER),
    ("Protects your branches", ["Warns you when a new site", "is too close to a store", "you already run."], RED),
]
cw = Inches(2.86)
gap = Inches(0.18)
x0 = Inches(0.85)
for i, (t, b, a) in enumerate(cards):
    card(s, x0 + i * (cw + gap), Inches(3.0), cw, Inches(2.5), t, b, accent=a,
         body_size=14)
text(s, Inches(0.85), Inches(5.8), Inches(11.6), Inches(0.9),
     [[("The big idea:  ", Pt(17), OAK_GREEN, True, False),
       ("instead of guessing, you make expansion decisions backed by data on demand, "
        "competition, rent, and your own branch network.", Pt(17), INK, False, False)]],
     line_spacing=1.1)
footer(s, page())  # page 2... but PAGE starts at 0; title is page 1 visually

# We'll handle numbering: reset so title isn't counted oddly.

# ===========================================================================
# SLIDE 3 — The 5 questions it answers
# ===========================================================================
s = add_slide(WHITE)
section_header(s, "THE CORE IDEA", "The 5 questions it answers for every site")
qs = [
    ("1.  Is there demand?", "Are there enough people and customers nearby to keep you busy?", LEAF),
    ("2.  Is the market crowded?", "How many competitors like you are already in the area? (We call open space “whitespace”.)", OAK_GREEN),
    ("3.  Does it fit my brand?", "Does the size, location and style suit how you operate — quick-service, dine-in, delivery or café?", AMBER),
    ("4.  Do the numbers work?", "Is the likely rent sensible compared to the sales the area can support?", LEAF),
    ("5.  Will it hurt my own stores?", "Is it far enough from your existing branches so they don't steal each other's customers?", RED),
]
y = Inches(1.7)
for title, body, accent in qs:
    rect(s, Inches(0.85), y, Inches(11.6), Inches(0.92), fill=SAND, radius=True)
    rect(s, Inches(0.85), y, Inches(0.12), Inches(0.92), fill=accent, radius=True)
    text(s, Inches(1.2), y + Inches(0.05), Inches(3.8), Inches(0.82),
         [P(title, 18, OAK_GREEN, bold=True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.0), y + Inches(0.05), Inches(7.2), Inches(0.82),
         [P(body, 15, INK)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y += Inches(1.03)
footer(s, page())

# ===========================================================================
# SLIDE 4 — How to run a search (step by step)
# ===========================================================================
s = add_slide(SAND)
section_header(s, "STEP 1 · TELL IT ABOUT YOUR BRAND", "Running a search in 5 easy steps")
steps = [
    ("Open the Expansion Brief form", "It's the panel on the left when you open the Expansion Advisor."),
    ("Enter your brand basics", "Brand name, food category (e.g. Burger, Coffee, Shawarma), and your service model: QSR, Dine-in, Delivery-first or Café."),
    ("Set the size you need", "A minimum and maximum space in square metres (for example 100–500 m²)."),
    ("Add your existing branches", "Optional but recommended — so it can keep new sites a safe distance away."),
    ("Press “Run Search”", "Sit back. In moments you get a ranked shortlist of the best sites."),
]
y = Inches(1.7)
for i, (head, body) in enumerate(steps, 1):
    circle = s.shapes.add_shape(9, Inches(0.95), y, Inches(0.7), Inches(0.7))  # 9 = oval
    circle.fill.solid(); circle.fill.fore_color.rgb = OAK_GREEN
    circle.line.fill.background(); circle.shadow.inherit = False
    tf = circle.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i); r.font.size = Pt(24); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Calibri"
    text(s, Inches(1.95), y - Inches(0.02), Inches(10.4), Inches(0.78),
         [[(head + "   ", Pt(18), OAK_GREEN, True, False),
           (body, Pt(15), SLATE, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y += Inches(1.0)
text(s, Inches(1.95), Inches(6.75), Inches(10.4), Inches(0.5),
     [P("Tip: in Advanced settings you can also set price tier, average check, and how much "
        "you care about parking, frontage and visibility.", 13, SLATE, italic=True)])
footer(s, page())

# ===========================================================================
# SCREENSHOT — The brief form
# ===========================================================================
screenshot_slide(
    "STEP 1 · THE REAL SCREEN", "This is the brief form you fill in",
    "ui-brief-form.png",
    "Everything the Advisor needs fits on one short form. Fill it in once and press the button at the bottom.",
    [
        ("Brand & category", "Who you are and what you sell.", LEAF),
        ("Service model", "QSR, dine-in, delivery-first or café.", OAK_GREEN),
        ("Area range", "The smallest and largest unit you'd take.", AMBER),
        ("Target districts", "Optional — pick areas or leave blank for all.", LEAF),
        ("Existing branches", "Add them so the Advisor protects your stores.", RED),
    ],
)

# ===========================================================================
# SLIDE 5 — What you get back (results overview)
# ===========================================================================
s = add_slide(WHITE)
section_header(s, "STEP 2 · READ YOUR RESULTS", "What the Advisor gives you back")
text(s, Inches(0.85), Inches(1.5), Inches(11.6), Inches(0.6),
     [P("You receive a ranked shortlist of up to 15 sites. Each one comes as a card you can tap to open.",
        17, INK)])

# Mock result card
cx, cy, cwid, chei = Inches(0.85), Inches(2.3), Inches(6.0), Inches(4.4)
rect(s, cx, cy, cwid, chei, fill=SAND, line=RGBColor(0xE2, 0xDD, 0xD2), radius=True)
# rank badge
badge = s.shapes.add_shape(9, cx + Inches(0.3), cy + Inches(0.3), Inches(0.85), Inches(0.85))
badge.fill.solid(); badge.fill.fore_color.rgb = OAK_GREEN; badge.line.fill.background()
badge.shadow.inherit = False
tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "#1"; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE
text(s, cx + Inches(1.35), cy + Inches(0.28), Inches(4.4), Inches(0.9),
     [P("Al Olaya · commercial unit", 18, INK, bold=True),
      P("Example candidate site", 12, SLATE)], space_after=Pt(2))
# score row
text(s, cx + Inches(0.35), cy + Inches(1.4), Inches(5.3), Inches(2.6),
     [[("Score  ", Pt(15), SLATE, False, False), ("87 / 100", Pt(15), OAK_GREEN, True, False)],
      [("Confidence  ", Pt(15), SLATE, False, False), ("Grade A", Pt(15), GREEN_OK, True, False)],
      [("Gate status  ", Pt(15), SLATE, False, False), ("Pass ✓", Pt(15), GREEN_OK, True, False)],
      [("Estimated rent  ", Pt(15), SLATE, False, False), ("~ SAR 480,000 / yr", Pt(15), INK, True, False)],
      [("Area  ", Pt(15), SLATE, False, False), ("220 m²", Pt(15), INK, True, False)],
      [("Tier  ", Pt(15), SLATE, False, False), ("Premier", Pt(15), LEAF, True, False)]],
     space_after=Pt(9))

# legend on right
ly = Inches(2.3)
legend = [
    ("Rank & Score", "Higher is better. 100 is a perfect match.", OAK_GREEN),
    ("Confidence (A–D)", "How sure we are. A = strong data, D = rough estimate.", LEAF),
    ("Gate status", "A quick safety check — Pass, Unknown, or Fail.", AMBER),
    ("Estimated rent", "Likely yearly rent for the unit.", INK),
    ("Tier", "Premier = best fit, Standard, or Exploratory.", RED),
]
rx = Inches(7.15)
for head, body, accent in legend:
    rect(s, rx, ly, Inches(5.3), Inches(0.78), fill=WHITE, line=RGBColor(0xE2, 0xDD, 0xD2), radius=True)
    rect(s, rx, ly, Inches(0.1), Inches(0.78), fill=accent, radius=True)
    text(s, rx + Inches(0.28), ly + Inches(0.04), Inches(5.0), Inches(0.72),
         [[(head + "  ", Pt(14), OAK_GREEN, True, False), (body, Pt(13), SLATE, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
    ly += Inches(0.86)
footer(s, page())

# ===========================================================================
# SCREENSHOT — The ranked results list
# ===========================================================================
screenshot_slide(
    "STEP 2 · THE REAL SCREEN", "Your ranked shortlist of sites",
    "ui-results.png",
    "Each site is a card, sorted best-first. The coloured tags give you the headline at a glance.",
    [
        ("Rank & score", "#1, #2, #3 … with a score out of 100.", OAK_GREEN),
        ("Lead Site / Premier", "Flags the standout, highest-quality sites.", LEAF),
        ("Pass / data grade", "Confidence grade and the safety-check result.", AMBER),
        ("Rent & size", "Annual rent, area and fit-out at a glance.", LEAF),
        ("+ / ! lines", "One key strength and one key risk per site.", RED),
    ],
)

# ===========================================================================
# SLIDE 6 — The traffic-light "gate" check
# ===========================================================================
s = add_slide(SAND)
section_header(s, "UNDERSTANDING RESULTS", "The traffic-light safety check (“Gates”)")
text(s, Inches(0.85), Inches(1.5), Inches(11.6), Inches(0.7),
     [P("Before a site reaches you, the Advisor runs quick safety checks — like zoning, distance "
        "from your branches, and whether the numbers add up. The result is a simple traffic light.",
        17, INK)], line_spacing=1.1)
lights = [
    ("PASS", GREEN_OK, "All checks cleared.", ["A low-risk site.", "Safe to shortlist and", "visit in person."]),
    ("UNKNOWN", AMBER, "Some data is missing.", ["No blocking problems,", "but do your own checks", "before committing."]),
    ("FAIL", RED, "A key check failed.", ["For example: wrong zoning,", "or too close to a branch", "you already run."]),
]
cw = Inches(3.7); x0 = Inches(0.85); gap = Inches(0.25)
for i, (label, col, sub, lines) in enumerate(lights):
    x = x0 + i * (cw + gap)
    rect(s, x, Inches(2.55), cw, Inches(3.6), fill=WHITE, line=RGBColor(0xE2, 0xDD, 0xD2), radius=True)
    dot = s.shapes.add_shape(9, x + cw/2 - Inches(0.55), Inches(2.85), Inches(1.1), Inches(1.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit = False
    text(s, x, Inches(4.15), cw, Inches(0.5), [P(label, 22, col, bold=True)], align=PP_ALIGN.CENTER)
    text(s, x, Inches(4.65), cw, Inches(0.4), [P(sub, 14, INK, bold=True)], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.3), Inches(5.1), cw - Inches(0.6), Inches(1.0),
         [P(l, 13, SLATE) for l in lines], align=PP_ALIGN.CENTER, space_after=Pt(2),
         line_spacing=1.0)
footer(s, page())

# ===========================================================================
# SLIDE 7 — Inside a site (detail + memo)
# ===========================================================================
s = add_slide(WHITE)
section_header(s, "STEP 3 · GO DEEPER", "Open a site to see the full story")
left = [
    ("Demand thesis", "Why customers are here — people nearby and competitor presence."),
    ("Cost thesis", "Likely rent and fit-out costs, and how they compare to similar sites."),
    ("Score breakdown", "How the score is built from demand, competition, fit and economics."),
    ("Top positives & risks", "The 2–3 best reasons to go — and the main things to watch."),
    ("Comparable competitors", "Nearby restaurants in your category, with their ratings."),
]
text(s, Inches(0.85), Inches(1.5), Inches(6.0), Inches(0.5),
     [P("Tap any card and the detail panel opens:", 17, INK, bold=True)])
y = Inches(2.1)
for head, body in left:
    rect(s, Inches(0.85), y, Inches(5.9), Inches(0.82), fill=SAND, radius=True)
    rect(s, Inches(0.85), y, Inches(0.1), Inches(0.82), fill=LEAF, radius=True)
    text(s, Inches(1.1), y + Inches(0.04), Inches(5.5), Inches(0.76),
         [[(head + "  ", Pt(15), OAK_GREEN, True, False), (body, Pt(13), SLATE, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
    y += Inches(0.9)

# Decision memo box
rect(s, Inches(7.1), Inches(2.1), Inches(5.35), Inches(4.5), fill=OAK_GREEN, radius=True)
text(s, Inches(7.45), Inches(2.3), Inches(4.7), Inches(0.6),
     [P("“View Decision Memo”", 19, WHITE, bold=True)])
text(s, Inches(7.45), Inches(2.9), Inches(4.7), Inches(3.5),
     [P("A short written recommendation, like an analyst wrote it for you:", 14, RGBColor(0xCF,0xE3,0xD7)),
      [("• Headline verdict  ", Pt(14), WHITE, True, False), ("— strong fit / consider carefully", Pt(13), RGBColor(0xCF,0xE3,0xD7), False, False)],
      [("• Key evidence  ", Pt(14), WHITE, True, False), ("— the main supporting facts", Pt(13), RGBColor(0xCF,0xE3,0xD7), False, False)],
      [("• Risks  ", Pt(14), WHITE, True, False), ("— what to watch out for", Pt(13), RGBColor(0xCF,0xE3,0xD7), False, False)],
      [("• Market research  ", Pt(14), WHITE, True, False), ("— delivery & district trends", Pt(13), RGBColor(0xCF,0xE3,0xD7), False, False)],
      [("• Best use case  ", Pt(14), WHITE, True, False), ("— how the site works best", Pt(13), RGBColor(0xCF,0xE3,0xD7), False, False)]],
     space_after=Pt(10), line_spacing=1.05)
text(s, Inches(7.45), Inches(6.05), Inches(4.7), Inches(0.5),
     [P("Perfect for sharing with partners or investors.", 13, RGBColor(0x9C,0xC4,0xAD), italic=True)])
footer(s, page())

# ===========================================================================
# SLIDE 8 — Comparing sites
# ===========================================================================
s = add_slide(SAND)
section_header(s, "STEP 4 · COMPARE", "Put your favourites head-to-head")
text(s, Inches(0.85), Inches(1.5), Inches(11.6), Inches(0.7),
     [P("Tick 2 to 6 sites and press Compare. You get a side-by-side table with a winner badge for "
        "each thing that matters — so the trade-offs are obvious at a glance.", 17, INK)],
     line_spacing=1.1)
badges = [
    "Best Overall", "Best Value", "Lowest Rent", "Best Economics",
    "Best Brand Fit", "Highest Demand", "Strongest Delivery Market",
    "Most Whitespace", "Lowest Cannibalization", "Most Confident",
]
# badge chips grid
bx, by = Inches(0.85), Inches(2.7)
chip_w, chip_h = Inches(2.78), Inches(0.7)
gx, gy = Inches(0.18), Inches(0.2)
for i, b in enumerate(badges):
    col = i % 4
    row = i // 4
    x = bx + col * (chip_w + gx)
    y = by + row * (chip_h + gy)
    rect(s, x, y, chip_w, chip_h, fill=WHITE, line=LEAF, line_w=Pt(1.5), radius=True)
    text(s, x, y, chip_w, chip_h, [[("★ ", Pt(14), AMBER, True, False),
                                     (b, Pt(13.5), OAK_GREEN, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(1.2),
     [[("Why this helps:  ", Pt(16), OAK_GREEN, True, False),
       ("rarely is one site best at everything. One may be cheapest, another may have the most "
        "customers. The badges let you pick the trade-off that fits your strategy.", Pt(16), INK, False, False)]],
     line_spacing=1.15)
footer(s, page())

# ===========================================================================
# SCREENSHOT — The comparison table
# ===========================================================================
screenshot_slide(
    "STEP 4 · THE REAL SCREEN", "Sites compared side-by-side",
    "ui-compare.png",
    "Pick 2–6 sites and the Advisor lays them out in one table. Green cells mark the winner on each row.",
    [
        ("Winner badges", "Best Overall, Best Value, Highest Demand…", OAK_GREEN),
        ("Grouped rows", "Demand, economics, rent and site quality.", LEAF),
        ("Green = best", "The strongest site per row is highlighted.", AMBER),
        ("Lead Site tag", "Marks the overall front-runner.", LEAF),
    ],
)

# ===========================================================================
# SLIDE 9 — Executive report / saving
# ===========================================================================
s = add_slide(WHITE)
section_header(s, "STEP 5 · DECIDE & SHARE", "Wrap it up into a report")
colA = [
    ("Executive Report", ["A clean one-screen summary:", "your top 3 sites, why the #1", "wins, the main risk, and the", "best store format."], OAK_GREEN),
    ("Presentation Mode", ["A tidy slideshow view, ready", "to show to partners,", "landlords or investors", "without any clutter."], LEAF),
]
colB = [
    ("Save your search", ["Keep a search as a Draft or", "mark it Final so you can", "return to it later."], AMBER),
    ("Shortlist finalists", ["Bookmark 2–6 favourite sites", "as finalists to revisit and", "compare another day."], RED),
]
y = Inches(1.7)
for head, lines, accent in colA:
    card(s, Inches(0.85), y, Inches(5.7), Inches(2.25), head, lines, accent=accent,
         body_size=14)
    y += Inches(2.5)
y = Inches(1.7)
for head, lines, accent in colB:
    card(s, Inches(6.75), y, Inches(5.7), Inches(2.25), head, lines, accent=accent,
         body_size=14)
    y += Inches(2.5)
footer(s, page())

# ===========================================================================
# SCREENSHOT — The executive report
# ===========================================================================
screenshot_slide(
    "STEP 5 · THE REAL SCREEN", "The one-page executive report",
    "ui-report.png",
    "When you're ready to decide, the report sums everything up on one screen — ideal to share with partners.",
    [
        ("Recommendation", "The headline call in one sentence.", OAK_GREEN),
        ("Why it wins / risk", "The single biggest plus and minus.", LEAF),
        ("Top candidates", "Your best 3 sites, side by side.", AMBER),
        ("Copy & Presentation", "Export or switch to a clean slideshow.", LEAF),
    ],
)

# ===========================================================================
# SLIDE 10 — Glossary
# ===========================================================================
s = add_slide(SAND)
section_header(s, "PLAIN-ENGLISH GLOSSARY", "The words you'll see, explained simply")
glossary = [
    ("Demand", "How many customers are likely nearby."),
    ("Whitespace", "Open opportunity — few competitors like you in the area."),
    ("Brand fit", "How well a site matches how your brand operates."),
    ("Economics", "Whether the likely rent makes sense for the sales possible."),
    ("Cannibalization", "Risk a new site steals sales from your own branches."),
    ("Confidence (A–D)", "How reliable the estimates are. A is best."),
    ("Gate", "A pass / fail safety check on a site."),
    ("Tier", "Quality label: Premier, Standard or Exploratory."),
    ("Delivery market", "How active food-delivery apps are in the area."),
    ("Provider density", "How many delivery platforms cover the area."),
    ("Comparable competitors", "Nearby rivals in your category, with ratings."),
    ("Shortlist", "Your saved set of promising sites to review."),
]
col_w = Inches(5.8)
x_cols = [Inches(0.85), Inches(6.75)]
y_start = Inches(1.65)
row_h = Inches(0.82)
for i, (term, mean) in enumerate(glossary):
    col = i // 6
    row = i % 6
    x = x_cols[col]
    y = y_start + row * row_h
    rect(s, x, y, col_w, Inches(0.72), fill=WHITE, line=RGBColor(0xE2,0xDD,0xD2), radius=True)
    rect(s, x, y, Inches(0.09), Inches(0.72), fill=LEAF, radius=True)
    text(s, x + Inches(0.25), y + Inches(0.03), col_w - Inches(0.4), Inches(0.66),
         [[(term + "  ", Pt(14.5), OAK_GREEN, True, False),
           (mean, Pt(12.5), SLATE, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.92)
footer(s, page())

# ===========================================================================
# SLIDE 11 — Quick tips / closing
# ===========================================================================
s = add_slide(OAK_GREEN)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=OAK_GREEN)
rect(s, 0, 0, Inches(0.35), SLIDE_H, fill=LEAF)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8),
     [P("A few tips to get the best results", 34, WHITE, bold=True)])
tips = [
    ("Always add your existing branches.", "It's the only way the Advisor can protect you from cannibalization."),
    ("Be honest about your service model.", "A delivery-first brand and a dine-in brand get scored very differently."),
    ("Start city-wide, then narrow.", "Leave target districts empty first to discover areas you hadn't considered."),
    ("Treat “Unknown” as “go check it”.", "It usually means missing data, not a bad site — a quick visit settles it."),
    ("Use the memo to win buy-in.", "It's written for partners and investors, not engineers."),
]
y = Inches(1.8)
for head, body in tips:
    chk = s.shapes.add_shape(9, Inches(0.95), y + Inches(0.05), Inches(0.4), Inches(0.4))
    chk.fill.solid(); chk.fill.fore_color.rgb = LEAF; chk.line.fill.background(); chk.shadow.inherit = False
    tf = chk.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "✓"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    text(s, Inches(1.6), y, Inches(10.8), Inches(0.85),
         [[(head + "   ", Pt(18), WHITE, True, False),
           (body, Pt(15), RGBColor(0xCF,0xE3,0xD7), False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y += Inches(0.92)
text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
     [P("Oaktree Atlas · Expansion Advisor — smarter F&B expansion in Riyadh.", 14,
        RGBColor(0x9C,0xC4,0xAD), italic=True)])

# ---------------------------------------------------------------------------
out = "docs/Expansion-Advisor-User-Guide.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides.")
